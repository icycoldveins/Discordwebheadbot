import discord
from discord.ext import commands
from discord import app_commands
import google.generativeai as genai
import os
from dotenv import load_dotenv
import random
import asyncio
import io
from pptx import Presentation
import docx
import PyPDF2
import json
import traceback
import re
import time
from typing import Optional
# 

class PresentationTriviaView(discord.ui.View):
    def __init__(self, cog, interaction):
        super().__init__(timeout=None)
        self.cog = cog
        self.interaction = interaction
        self.active = True
        self.answer_event = asyncio.Event()
        self.last_reaction = None
        self.last_user = None

    @discord.ui.button(label="End Trivia", style=discord.ButtonStyle.danger)
    async def end_trivia(self, interaction: discord.Interaction, button: discord.ui.Button):
        if interaction.user.id == self.interaction.user.id:
            self.active = False
            self.answer_event.set()
            for item in self.children:
                item.disabled = True
            await interaction.message.edit(view=self)
            
            # Clear user's cache when they end trivia
            self.cog.clear_user_cache(interaction.user.id)
            
            await interaction.response.send_message("Trivia ended by user.")
        else:
            await interaction.response.send_message("Only the person who started the trivia can end it.", ephemeral=True)

    async def wait_for_answer(self, bot, message, valid_reactions):
        try:
            def check(reaction, user):
                return (
                    user.id == self.interaction.user.id 
                    and str(reaction.emoji) in valid_reactions 
                    and reaction.message.id == message.id
                )

            reaction, user = await bot.wait_for('reaction_add', timeout=30.0, check=check)
            self.last_reaction = reaction
            self.last_user = user
            self.answer_event.set()
            return reaction, user

        except asyncio.TimeoutError:
            self.answer_event.set()
            return None, None

class ContentChoiceView(discord.ui.View):
    def __init__(self, cog, interaction):
        super().__init__(timeout=300)
        self.cog = cog
        self.interaction = interaction

    @discord.ui.button(label="Use Previous Content", style=discord.ButtonStyle.green, emoji="🔄")
    async def use_previous(self, interaction: discord.Interaction, button: discord.ui.Button):
        await interaction.response.defer()
        content = self.cog.user_content.get(interaction.user.id)
        if content:
            # Disable buttons after selection
            for item in self.children:
                item.disabled = True
            await interaction.message.edit(view=self)
            
            processing_msg = await interaction.followup.send("Generating quiz from previous content... 🎯")
            await self.cog.start_quiz(interaction, content)
        else:
            await interaction.followup.send("No previous content found. Please provide new content.")

    @discord.ui.button(label="New Content", style=discord.ButtonStyle.primary, emoji="📝")
    async def new_content(self, interaction: discord.Interaction, button: discord.ui.Button):
        # Disable buttons after selection
        for item in self.children:
            item.disabled = True
        await interaction.message.edit(view=self)
        
        await interaction.response.send_message(
            "Please paste your content or upload a file (txt, pdf, docx, pptx). You have 5 minutes."
        )

        def check(m):
            return (m.author == interaction.user and 
                   m.channel == interaction.channel and 
                   (m.content or m.attachments))

        try:
            msg = await self.cog.bot.wait_for('message', timeout=300.0, check=check)
            
            processing_msg = await interaction.followup.send("Processing your content... ⏳")
            
            if msg.attachments:
                file = msg.attachments[0]
                if not any(file.filename.endswith(ext) for ext in ['.txt', '.pdf', '.docx', '.pptx']):
                    await interaction.followup.send(
                        "Invalid file type. Please use .txt, .pdf, .docx, or .pptx files."
                    )
                    return
                    
                file_data = await file.read()
                content = await self.cog.process_file(file_data, file.filename, interaction)
            else:
                content = msg.content

            if not content or len(content.strip()) < 50:
                await interaction.followup.send(
                    "Not enough content provided. Please try again with more text."
                )
                return
            
            await processing_msg.edit(content="Content processed! Generating quiz... 🎯")
            await self.cog.start_quiz(interaction, content)
                
        except asyncio.TimeoutError:
            await interaction.followup.send("No content received within 5 minutes. Please try again.")

class PresentationTrivia(commands.Cog):
    def __init__(self, bot):
        self.bot = bot
        # Initialize Gemini
        load_dotenv()
        genai.configure(api_key=os.getenv('GEMINI_API_KEY'))
        self.model = genai.GenerativeModel('gemini-2.0-flash')
        
        # Simplified cache structure
        self.user_content = {}          # {user_id: str}  # Just store the content string
        self.chunk_cache = {}           # Quiz-specific cache
        self.used_questions = {}        # Quiz-specific cache
        self.chunk_size = 4000          # Size of content chunk to process at a time
        self.MAX_CONTENT_SIZE = 100000  # Maximum content size per user
        
    async def generate_questions(self, content: str, start_pos: int = 0, user_id: int = None, interaction: discord.Interaction = None) -> list:
        """Generate questions using Gemini AI from a random unused chunk of content"""
        try:
            if not content or len(content.strip()) < 50:
                print("Content too short or empty")
                return []

            # Calculate total chunks
            total_chunks = (len(content) + self.chunk_size - 1) // self.chunk_size
            
            # Initialize or get the remaining chunks list for this user
            if user_id not in self.chunk_cache:
                # Create list of all chunk indices
                self.chunk_cache[user_id] = {'chunks': list(range(total_chunks)), 'total_chunks': total_chunks, 'processed': {}, 'last_access': time.time()}
                # Shuffle the list initially
                random.shuffle(self.chunk_cache[user_id]['chunks'])
            
            # If no chunks left, we're done
            if not self.chunk_cache[user_id]['chunks']:
                print("All chunks have been used")
                return []

            # Pop the next chunk index from our shuffled list
            chosen_chunk = self.chunk_cache[user_id]['chunks'].pop()
            chunk_start = chosen_chunk * self.chunk_size
            
            # Extract the chunk
            content_chunk = content[chunk_start:chunk_start + self.chunk_size]
            print(f"Processing chunk {chosen_chunk + 1}/{total_chunks}, length: {len(content_chunk)} characters")

            # Send processing message
            processing_embed = discord.Embed(
                title="🤖 Processing Content",
                description=f"Processing chunk {chosen_chunk + 1}/{total_chunks}...\n" +
                           f"Chunks remaining: {len(self.chunk_cache[user_id]['chunks'])}/{total_chunks}",
                color=discord.Color.blue()
            )
            processing_msg = await interaction.followup.send(embed=processing_embed)

            prompt = f"""You are a quiz generator. Generate multiple-choice questions about the key concepts from this content chunk.

REQUIRED FORMAT (with proper spacing):
{{
    "questions": [
        {{
            "question": "What is the main concept described in this passage?",
            "correct_answer": "This is the correct answer with proper spacing",
            "incorrect_answers": [
                "First wrong answer with proper spacing",
                "Second wrong answer with proper spacing",
                "Third wrong answer with proper spacing"
            ],
            "explanation": "This is the explanation with proper spacing"
        }}
    ]
}}

CRITICAL RULES:
1. Generate as many questions as the content allows (no limit)
2. Only create questions when there's enough content to support them
3. Focus on key concepts and important information
4. NO questions about metadata (instructors, chapters, etc.)
5. Each question must have clear, distinct answers
6. All answers must come from the content
7. Quality over quantity
8. IMPORTANT: Maintain proper spacing between words
9. Use complete sentences
10. Add punctuation with spaces

Content chunk to use:
{content_chunk}"""

            print("Sending request to Gemini...")
            response = await asyncio.to_thread(
                self.model.generate_content, prompt
            )
            
            def clean_text(text):
                # Add space after punctuation
                text = re.sub(r'([.!?,])([A-Za-z])', r'\1 \2', text)
                # Add space between words if missing
                text = re.sub(r'([a-z])([A-Z])', r'\1 \2', text)
                # Fix multiple spaces
                text = re.sub(r'\s+', ' ', text)
                # Ensure proper spacing around quotes
                text = re.sub(r'"(\w)', r'" \1', text)
                text = re.sub(r'(\w)"', r'\1 "', text)
                return text.strip()

            try:
                json_match = re.search(r'\{.*\}', response.text, re.DOTALL)
                if not json_match:
                    return []
                
                parsed = json.loads(json_match.group())
                if 'questions' not in parsed:
                    return []
                
                metadata_keywords = ['instructor', 'professor', 'teacher', 'chapter', 
                                   'section', 'page', 'course', 'code', 'syllabus']
                
                valid_questions = []
                for q in parsed['questions']:
                    # Create a unique hash for this question
                    question_hash = hash(f"{q['question']}:{q['correct_answer']}")
                    
                    # Skip if we've used this question before
                    if question_hash in self.used_questions[user_id]:
                        continue

                    if any(keyword in q['question'].lower() for keyword in metadata_keywords):
                        continue
                        
                    if (isinstance(q, dict) and
                        all(key in q for key in ['question', 'correct_answer', 'incorrect_answers', 'explanation']) and
                        isinstance(q['incorrect_answers'], list) and 
                        len(q['incorrect_answers']) == 3):
                        
                        # Clean and format all text fields
                        cleaned_q = {
                            'question': clean_text(q['question']),
                            'correct_answer': clean_text(q['correct_answer']),
                            'incorrect_answers': [clean_text(ans) for ans in q['incorrect_answers']],
                            'explanation': clean_text(q['explanation'])
                        }
                        
                        # Validate answer uniqueness and quality
                        answers = [cleaned_q['correct_answer']] + cleaned_q['incorrect_answers']
                        if (len(set(answers)) == 4 and  # All answers must be unique
                            all(len(ans.strip()) > 0 for ans in answers)):  # No empty answers
                            # Add question hash to used set
                            self.used_questions[user_id].add(question_hash)
                            valid_questions.append(cleaned_q)
                
                if valid_questions:
                    random.shuffle(valid_questions)
                    self.chunk_cache[user_id]['processed'][chosen_chunk] = valid_questions
                    return valid_questions
                
            except Exception as e:
                print(f"Error processing questions: {str(e)}")
                return []
                
        except Exception as e:
            print(f"Question generation error: {str(e)}")
            return []

    def clean_content(self, content: str) -> str:
        """Clean and prepare content for question generation"""
        if not content:
            return ""
            
        # Remove extra whitespace and normalize line endings
        content = content.strip()
        content = ' '.join(content.split())
        
        # Remove any empty lines
        content = '\n'.join(line.strip() for line in content.split('\n') if line.strip())
        
        # Remove any special characters that might interfere with processing
        content = content.replace('\x00', '')
        
        return content

    async def process_file(self, file_data: bytes, filename: str, interaction: discord.Interaction = None) -> str:
        """Process different file types and extract text contents"""
        try:
            text_content = ""  # Initialize text_content variable
            
            if filename.endswith('.pptx'):
                presentation = Presentation(io.BytesIO(file_data))
                text_content = []
                
                for slide in presentation.slides:
                    # Extract text from shapes
                    for shape in slide.shapes:
                        if hasattr(shape, "text") and shape.text.strip():
                            text_content.append(shape.text.strip())
                        
                        # Handle tables specifically
                        if shape.has_table:
                            for row in shape.table.rows:
                                for cell in row.cells:
                                    if cell.text.strip():
                                        text_content.append(cell.text.strip())
                        
                        # Handle text frames
                        if hasattr(shape, "text_frame"):
                            for paragraph in shape.text_frame.paragraphs:
                                if paragraph.text.strip():
                                    text_content.append(paragraph.text.strip())

                return "\n".join(text_content)
            
            elif filename.endswith('.docx'):
                doc = docx.Document(io.BytesIO(file_data))
                for paragraph in doc.paragraphs:
                    if paragraph.text.strip():
                        text_content += paragraph.text.strip() + "\n"
                    
            elif filename.endswith('.pdf'):
                pdf = PyPDF2.PdfReader(io.BytesIO(file_data))
                for page in pdf.pages:
                    text = page.extract_text()
                    if text.strip():
                        text_content += text.strip() + "\n"
                    
            elif filename.endswith('.txt'):
                text_content = file_data.decode('utf-8')
            
            # Clean and validate content
            cleaned_content = self.clean_content(text_content)
            if len(cleaned_content) < 50:  # Minimum content length check
                print(f"Content too short after cleaning: {len(cleaned_content)} chars")
                return ""
                
            print(f"Successfully processed {filename}, extracted {len(cleaned_content)} chars")
            return cleaned_content
            
        except Exception as e:
            print(f"Error processing file {filename}: {str(e)}")
            if interaction:  # Only send message if interaction is provided
                await interaction.followup.send(
                    f"Error processing file: {str(e)}. Please try a different file or paste the content directly."
                )
            return ""

    @app_commands.command(name="quiz", description="Start a quiz game from presentation content")
    @app_commands.describe(new_content="Start with new content? Default: Use previous content if available")
    async def presentation_trivia(self, interaction: discord.Interaction, new_content: bool = False):
        # Defer the response immediately to prevent timeout
        await interaction.response.defer()

        if not new_content and interaction.user.id in self.user_content:
            embed = discord.Embed(
                title="📚 Presentation Quiz",
                description="Would you like to use your previous content or provide new content?",
                color=discord.Color.blue()
            )
            view = ContentChoiceView(self, interaction)
            await interaction.followup.send(embed=embed, view=view)
            return
            
        # If no previous content or user wants new content
        await interaction.followup.send(
            "Please paste your content or upload a file (txt, pdf, docx, pptx). You have 5 minutes."
        )

        def check(m):
            return (m.author == interaction.user and 
                   m.channel == interaction.channel and 
                   (m.content or m.attachments))

        try:
            msg = await self.bot.wait_for('message', timeout=300.0, check=check)
            
            processing_msg = await interaction.followup.send("Processing your content... ⏳")
            
            if msg.attachments:
                file = msg.attachments[0]
                if not any(file.filename.endswith(ext) for ext in ['.txt', '.pdf', '.docx', '.pptx']):
                    await interaction.followup.send(
                        "Invalid file type. Please use .txt, .pdf, .docx, or .pptx files."
                    )
                    return
                    
                file_data = await file.read()
                content = await self.process_file(file_data, file.filename, interaction)
            else:
                content = msg.content

            if not content or len(content.strip()) < 50:
                await interaction.followup.send(
                    "Not enough content provided. Please try again with more text."
                )
                return
            
            await processing_msg.edit(content="Content processed! Generating quiz... 🎯")
            await self.start_quiz(interaction, content)
                
        except asyncio.TimeoutError:
            await interaction.followup.send("No content received within 5 minutes. Please try again.")

    def clear_user_cache(self, user_id: int):
        """Clear quiz-specific cache for a user"""
        self.chunk_cache.pop(user_id, None)
        self.used_questions.pop(user_id, None)
        # Note: We don't clear user_content here to preserve it for future use

    async def start_quiz(self, interaction: discord.Interaction, content: str):
        """Start the quiz with the provided content"""
        # Clear quiz-specific cache
        self.clear_user_cache(interaction.user.id)
        
        # Check content size
        if len(content) > self.MAX_CONTENT_SIZE:
            await interaction.followup.send(
                f"Content too large. Please provide content under {self.MAX_CONTENT_SIZE/1000}KB."
            )
            return
            
        # Store the content for future use
        self.user_content[interaction.user.id] = content
        
        # Initialize quiz-specific tracking data
        self.used_questions[interaction.user.id] = set()
        
        try:
            view = PresentationTriviaView(self, interaction)
            score = 0
            current_question = 0
            content_position = 0
            questions = []

            # Initial progress message
            progress = discord.Embed(
                title="📚 Presentation Trivia",
                description=f"Current Score: {score}/{current_question}\nGenerating questions...",
                color=discord.Color.blue()
            )
            await interaction.followup.send(embed=progress)

            while view.active:
                if not questions:
                    # Check if trivia was ended before generating new questions
                    if not view.active:
                        break

                    questions = await self.generate_questions(
                        content, 
                        content_position,
                        interaction.user.id,
                        interaction
                    )
                    
                    if not questions:
                        # Add this condition to exit when no more questions can be generated
                        final_embed = discord.Embed(
                            title="🎯 Quiz Complete!",
                            description=f"All content has been covered!\nFinal Score: {score}/{current_question}",
                            color=discord.Color.gold()
                        )
                        await interaction.followup.send(embed=final_embed)
                        break

                # Get current question
                question = questions.pop(0)  # Take the next question
                
                # Format question
                question_embed = discord.Embed(
                    title=f"Question {current_question + 1}",
                    description=question['question'],
                    color=discord.Color.blue()
                )

                all_answers = question['incorrect_answers'] + [question['correct_answer']]
                random.shuffle(all_answers)

                for idx, answer in enumerate(all_answers, 1):
                    question_embed.add_field(
                        name=f"Option {idx}",
                        value=answer,
                        inline=False
                    )

                msg = await interaction.followup.send(embed=question_embed, view=view)

                # Add reactions
                reactions = ["1️⃣", "2️⃣", "3️⃣", "4️⃣"]
                for reaction in reactions:
                    await msg.add_reaction(reaction)

                # Check if trivia was ended before waiting for answer
                if not view.active:
                    break

                reaction, user = await view.wait_for_answer(self.bot, msg, reactions)
                
                if not view.active:  # If trivia was ended
                    final_embed = discord.Embed(
                        title="🎯 Quiz Complete!",
                        description=f"All content has been covered!\nFinal Score: {score}/{current_question}",
                        color=discord.Color.gold()
                    )
                    await interaction.followup.send(embed=final_embed)
                    break
                    
                if not reaction:  # If we timed out
                    timeout_embed = discord.Embed(
                        title="⏰ Time's Up!",
                        description=f"The correct answer was: **{question['correct_answer']}**\nFinal Score: {score}/{current_question}",
                        color=discord.Color.red()
                    )
                    await interaction.followup.send(embed=timeout_embed)
                    break

                # Process answer
                selected_idx = reactions.index(str(reaction.emoji))
                selected_answer = all_answers[selected_idx]
                is_correct = selected_answer == question['correct_answer']

                if is_correct:
                    score += 1

                # Show answer result
                result_embed = discord.Embed(
                    title="✨ Answer Result",
                    description=(
                        f"{'✅ Correct!' if is_correct else '❌ Incorrect!'}\n"
                        f"You selected: {selected_answer}\n"
                        f"Correct answer: **{question['correct_answer']}**\n\n"
                        f"Explanation: {question.get('explanation', 'No explanation provided.')}\n\n"
                        f"Current Score: {score}/{current_question + 1}"
                    ),
                    color=discord.Color.green() if is_correct else discord.Color.red()
                )
                await interaction.followup.send(embed=result_embed)
                
                # Move to next question
                current_question += 1
                
                # Add a small delay between questions
                if view.active:
                    await asyncio.sleep(2)

            # Show final score if we completed all questions
            if view.active and current_question == len(questions):
                final_embed = discord.Embed(
                    title="🎉 Quiz Complete!",
                    description=f"Congratulations! You've completed all questions!\nFinal Score: {score}/{current_question}\nPercentage: {(score/current_question)*100:.1f}%",
                    color=discord.Color.gold()
                )
                await interaction.followup.send(embed=final_embed)

        except Exception as e:
            if view.active:  # Only show error if trivia wasn't manually ended
                traceback.print_exc()
                await interaction.followup.send(f"An error occurred: {str(e)}")

    async def get_next_question(self, user_id: int) -> Optional[dict]:
        content = self.user_content.get(user_id)
        if not content:
            return None

        # Get or create chunks for this content
        chunks = self.chunk_cache.get(user_id)
        if not chunks:
            chunks = [content[i:i + self.chunk_size] for i in range(0, len(content), self.chunk_size)]
            self.chunk_cache[user_id] = chunks

        used_questions = self.used_questions.get(user_id, set())
        
        # Try each chunk until we get a valid question
        if chunks:  # Only if we have chunks left
            # Swap random chunk with last chunk and pop
            if len(chunks) > 1:  # Only need to swap if more than 1 chunk
                rand_idx = random.randint(0, len(chunks) - 2)  # -2 to avoid picking last index
                chunks[rand_idx], chunks[-1] = chunks[-1], chunks[rand_idx]
            chunk = chunks.pop()  # Always pop from end
            
            try:
                # Generate questions from this chunk
                questions = await self.generate_questions(chunk)
                
                # Find unused questions
                available_questions = [q for q in questions if q['id'] not in used_questions]
                if available_questions:
                    # Same swap-and-pop for question selection
                    if len(available_questions) > 1:
                        rand_idx = random.randint(0, len(available_questions) - 2)
                        available_questions[rand_idx], available_questions[-1] = available_questions[-1], available_questions[rand_idx]
                    question = available_questions.pop()
                    used_questions.add(question['id'])
                    self.used_questions[user_id] = used_questions
                    return question
                    
            except Exception as e:
                print(f"Error generating question: {e}")

async def setup(bot):
    await bot.add_cog(PresentationTrivia(bot)) 