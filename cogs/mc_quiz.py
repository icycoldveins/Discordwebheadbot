import discord
from discord import app_commands
from discord.ext import commands
import asyncio
import io
from pptx import Presentation
import docx
import PyPDF2
import os
from dotenv import load_dotenv
import google.generativeai as genai
import traceback
from .presentation_trivia import PresentationTriviaView
import random
import json

class MCQuestionView(discord.ui.View):
    def __init__(self, quiz_session):
        super().__init__()
        self.quiz_session = quiz_session

    @discord.ui.button(label='A', style=discord.ButtonStyle.primary)
    async def button_a(self, interaction: discord.Interaction, button: discord.ui.Button):
        await interaction.response.defer()
        await self.quiz_session.handle_answer('A')
        self.disable_all_buttons()

    @discord.ui.button(label='B', style=discord.ButtonStyle.primary)
    async def button_b(self, interaction: discord.Interaction, button: discord.ui.Button):
        await interaction.response.defer()
        await self.quiz_session.handle_answer('B')
        self.disable_all_buttons()

    @discord.ui.button(label='C', style=discord.ButtonStyle.primary)
    async def button_c(self, interaction: discord.Interaction, button: discord.ui.Button):
        await interaction.response.defer()
        await self.quiz_session.handle_answer('C')
        self.disable_all_buttons()

    @discord.ui.button(label='D', style=discord.ButtonStyle.primary)
    async def button_d(self, interaction: discord.Interaction, button: discord.ui.Button):
        await interaction.response.defer()
        await self.quiz_session.handle_answer('D')
        self.disable_all_buttons()

    def disable_all_buttons(self):
        for item in self.children:
            item.disabled = True

class ContentChoiceView(discord.ui.View):
    def __init__(self, cog, interaction):
        super().__init__()
        self.cog = cog
        self.interaction = interaction

    @discord.ui.button(label='Use Previous Content', style=discord.ButtonStyle.primary)
    async def use_previous(self, interaction: discord.Interaction, button: discord.ui.Button):
        await interaction.response.defer()
        content = self.cog.user_content.get(interaction.user.id, "")
        await self.cog.start_mc_quiz(interaction, content)

    @discord.ui.button(label='Provide New Content', style=discord.ButtonStyle.secondary)
    async def new_content(self, interaction: discord.Interaction, button: discord.ui.Button):
        await interaction.response.defer()
        await self.cog.mc_quiz(interaction, True)

class MCQuiz(commands.Cog):
    def __init__(self, bot):
        self.bot = bot
        # Initialize Gemini like PresentationTrivia
        load_dotenv()
        genai.configure(api_key=os.getenv('GEMINI_API_KEY'))
        self.model = genai.GenerativeModel('gemini-pro')
        self.user_content = {}
        self.used_questions = {}  # Track used questions per user

    @app_commands.command(name="mcquiz", description="Generate multiple-choice questions on a topic or from content")
    @app_commands.describe(new_content="Start with new content? Default: Use previous content if available")
    async def mc_quiz(self, interaction: discord.Interaction, new_content: bool = False):
        await interaction.response.defer()

        if not new_content and interaction.user.id in self.user_content:
            embed = discord.Embed(
                title="📝 Multiple Choice Quiz",
                description="Would you like to use your previous content or provide new content/topic?",
                color=discord.Color.blue()
            )
            view = ContentChoiceView(self, interaction)
            await interaction.followup.send(embed=embed, view=view)
            return

        await interaction.followup.send(
            "Please provide a topic (e.g., 'Python Programming Basics') or paste your content. "
            "You can also upload a file (txt, pdf, docx, pptx). You have 5 minutes."
        )

        def check(m):
            return (m.author == interaction.user and 
                   m.channel == interaction.channel and 
                   (m.content or m.attachments))

        try:
            msg = await self.bot.wait_for('message', timeout=300.0, check=check)
            
            processing_msg = await interaction.followup.send("Processing your input... ⏳")
            
            if msg.attachments:
                file = msg.attachments[0]
                if not any(file.filename.endswith(ext) for ext in ['.txt', '.pdf', '.docx', '.pptx']):
                    await interaction.followup.send(
                        "Invalid file type. Please use .txt, .pdf, .docx, or .pptx files."
                    )
                    return
                    
                file_data = await file.read()
                content = await self.process_file(file_data, file.filename, interaction)
            else:
                content = msg.content

            if not content or len(content.strip()) < 3:
                await interaction.followup.send(
                    "Please provide a valid topic or content."
                )
                return
            
            await processing_msg.edit(content="Input processed! Generating multiple choice quiz... 🎯")
            await self.start_mc_quiz(interaction, content)
                
        except asyncio.TimeoutError:
            await interaction.followup.send("No input received within 5 minutes. Please try again.")

    async def process_file(self, file_data: bytes, filename: str, interaction: discord.Interaction = None) -> str:
        """Process different file types and extract text contents"""
        try:
            text_content = ""
            
            if filename.endswith('.pptx'):
                presentation = Presentation(io.BytesIO(file_data))
                text_content = []
                
                for slide in presentation.slides:
                    for shape in slide.shapes:
                        if hasattr(shape, "text") and shape.text.strip():
                            text_content.append(shape.text.strip())
                        
                        if shape.has_table:
                            for row in shape.table.rows:
                                for cell in row.cells:
                                    if cell.text.strip():
                                        text_content.append(cell.text.strip())
                        
                        if hasattr(shape, "text_frame"):
                            for paragraph in shape.text_frame.paragraphs:
                                if paragraph.text.strip():
                                    text_content.append(paragraph.text.strip())

                return "\n".join(text_content)
            
            elif filename.endswith('.docx'):
                doc = docx.Document(io.BytesIO(file_data))
                for paragraph in doc.paragraphs:
                    if paragraph.text.strip():
                        text_content += paragraph.text.strip() + "\n"
                    
            elif filename.endswith('.pdf'):
                pdf = PyPDF2.PdfReader(io.BytesIO(file_data))
                for page in pdf.pages:
                    text = page.extract_text()
                    if text.strip():
                        text_content += text.strip() + "\n"
                    
            elif filename.endswith('.txt'):
                text_content = file_data.decode('utf-8')
            
            return text_content.strip()
            
        except Exception as e:
            print(f"Error processing file {filename}: {str(e)}")
            if interaction:
                await interaction.followup.send(
                    f"Error processing file: {str(e)}. Please try a different file or paste the content directly."
                )
            return ""

    async def generate_questions(self, topic: str, user_id: int = None, interaction: discord.Interaction = None) -> list:
        """Generate multiple choice questions about a topic using Gemini AI"""
        try:
            if user_id not in self.used_questions:
                self.used_questions[user_id] = set()

            prompt = f"""You are a quiz generator. Generate multiple-choice questions about this topic.

REQUIRED FORMAT (with proper spacing):
{{
    "questions": [
        {{
            "question": "Write a clear, specific question about {topic}?",
            "correct_answer": "The correct answer with proper spacing",
            "incorrect_answers": [
                "First wrong answer with proper spacing",
                "Second wrong answer with proper spacing",
                "Third wrong answer with proper spacing"
            ],
            "explanation": "Brief explanation of why the correct answer is right"
        }}
    ]
}}

CRITICAL RULES:
1. Generate 15 questions about the topic
2. Each question must test different aspects and concepts
3. Questions should have a mix of:
   - Basic knowledge questions
   - Intermediate understanding questions
   - Advanced application questions
4. Each question must have clear, distinct answers
5. All answers must be factually correct
6. Quality over quantity
7. IMPORTANT: Maintain proper spacing between words
8. Use complete sentences
9. Add punctuation with spaces
10. Include brief but informative explanations

Topic to use: {topic}"""

            print("Sending request to Gemini...")
            response = await asyncio.to_thread(
                self.model.generate_content, prompt
            )
            
            try:
                # Clean the response text
                response_text = response.text.strip()
                if response_text.startswith('```json'):
                    response_text = response_text[7:-3]  # Remove ```json and ``` markers
                
                # Parse the JSON response
                questions_data = json.loads(response_text)
                questions = questions_data.get('questions', [])
                
                if not questions:
                    print("No questions found in response")
                    return []

                # Filter out previously used questions
                questions = [q for q in questions 
                           if hash(q.get('question', '')) not in self.used_questions[user_id]]
                
                # Validate question format
                valid_questions = []
                for q in questions:
                    if all(key in q for key in ['question', 'correct_answer', 'incorrect_answers', 'explanation']):
                        if isinstance(q['incorrect_answers'], list) and len(q['incorrect_answers']) == 3:
                            valid_questions.append(q)
                
                # Track newly used questions
                for q in valid_questions:
                    self.used_questions[user_id].add(hash(q['question']))
                
                print(f"Generated {len(valid_questions)} valid questions")
                return valid_questions

            except json.JSONDecodeError as e:
                print(f"Error parsing JSON: {str(e)}")
                print(f"Raw response: {response.text}")
                
                # Attempt to clean and retry parsing
                try:
                    # Remove any potential markdown formatting
                    cleaned_text = response.text.replace('```json\n', '').replace('\n```', '').strip()
                    questions_data = json.loads(cleaned_text)
                    questions = questions_data.get('questions', [])
                    return questions
                except:
                    print("Failed to parse even after cleaning")
                    return []

        except Exception as e:
            print(f"Error generating questions: {str(e)}")
            if interaction:
                await interaction.followup.send("Error generating questions. Please try again.")
            return []

    async def start_mc_quiz(self, interaction: discord.Interaction, topic: str):
        """Start a continuous multiple choice quiz session on a topic"""
        try:
            self.user_content[interaction.user.id] = topic
            view = PresentationTriviaView(self, interaction)
            score = 0
            current_question = 0

            # Initialize used questions set if it doesn't exist
            if interaction.user.id not in self.used_questions:
                self.used_questions[interaction.user.id] = set()

            progress = discord.Embed(
                title="📝 Multiple Choice Quiz",
                description=f"Topic: {topic}\nCurrent Score: {score}/{current_question}\nGenerating questions...\nPreviously answered questions: {len(self.used_questions[interaction.user.id])}",
                color=discord.Color.blue()
            )
            await interaction.followup.send(embed=progress)

            while view.active:
                # Generate new batch of questions (without clearing used_questions)
                questions = await self.generate_questions(topic, interaction.user.id, interaction)
                
                if not questions:
                    final_embed = discord.Embed(
                        title="⚠️ Generation Error",
                        description=f"Failed to generate new questions. {len(self.used_questions[interaction.user.id])} questions have been asked about this topic.\nTry a different topic or aspect!",
                        color=discord.Color.red()
                    )
                    await interaction.followup.send(embed=final_embed)
                    break

                # Notify user of new batch
                batch_embed = discord.Embed(
                    title="🎯 New Questions Generated!",
                    description=f"Generated {len(questions)} new questions about {topic}.\nTotal unique questions asked: {len(self.used_questions[interaction.user.id])}",
                    color=discord.Color.green()
                )
                await interaction.followup.send(embed=batch_embed)
                await asyncio.sleep(2)

                for question in questions:
                    if not view.active:
                        break

                    # Format question
                    question_embed = discord.Embed(
                        title=f"Question {current_question + 1}",
                        description=question['question'],
                        color=discord.Color.blue()
                    )

                    all_answers = question['incorrect_answers'] + [question['correct_answer']]
                    random.shuffle(all_answers)

                    for idx, answer in enumerate(all_answers):
                        question_embed.add_field(
                            name=f"{chr(65 + idx)})",  # A), B), C), D)
                            value=answer,
                            inline=False
                        )

                    msg = await interaction.followup.send(embed=question_embed, view=view)

                    # Add reactions
                    reactions = ["🇦", "🇧", "🇨", "🇩"]
                    for reaction in reactions:
                        await msg.add_reaction(reaction)

                    reaction, user = await view.wait_for_answer(self.bot, msg, reactions)
                    
                    if not view.active:  # If quiz was ended
                        final_embed = discord.Embed(
                            title="🎯 Quiz Ended!",
                            description=f"Final Score: {score}/{current_question}\nThanks for playing!",
                            color=discord.Color.gold()
                        )
                        await interaction.followup.send(embed=final_embed)
                        return

                    if not reaction:  # If we timed out
                        timeout_embed = discord.Embed(
                            title="⏰ Time's Up!",
                            description=f"The correct answer was: **{question['correct_answer']}**\nFinal Score: {score}/{current_question}",
                            color=discord.Color.red()
                        )
                        await interaction.followup.send(embed=timeout_embed)
                        break

                    # Process answer
                    selected_idx = reactions.index(str(reaction.emoji))
                    selected_answer = all_answers[selected_idx]
                    is_correct = selected_answer == question['correct_answer']

                    if is_correct:
                        score += 1

                    # Show answer result
                    result_embed = discord.Embed(
                        title="✨ Answer Result",
                        description=(
                            f"{'✅ Correct!' if is_correct else '❌ Incorrect!'}\n"
                            f"You selected: {selected_answer}\n"
                            f"Correct answer: **{question['correct_answer']}**\n\n"
                            f"Explanation: {question.get('explanation', 'No explanation provided.')}\n\n"
                            f"Current Score: {score}/{current_question + 1}"
                        ),
                        color=discord.Color.green() if is_correct else discord.Color.red()
                    )
                    await interaction.followup.send(embed=result_embed)
                    
                    current_question += 1
                    
                    if view.active:
                        await asyncio.sleep(2)

        except Exception as e:
            if view.active:
                traceback.print_exc()
                await interaction.followup.send(f"An error occurred: {str(e)}")

class MCQuizSession:
    def __init__(self, bot, interaction, questions):
        self.bot = bot
        self.interaction = interaction
        self.questions = questions
        self.current_question = 0
        self.score = 0
        self.total_questions = len(questions)

    async def start(self):
        """Start the quiz session"""
        await self.show_question()

    async def show_question(self):
        """Display the current question"""
        if self.current_question >= len(self.questions):
            await self.end_quiz()
            return

        question = self.questions[self.current_question]
        
        embed = discord.Embed(
            title=f"Question {self.current_question + 1}/{self.total_questions}",
            description=question['question'],
            color=discord.Color.blue()
        )

        for i, option in enumerate(question['options']):
            embed.add_field(
                name=f"{chr(65 + i)})",
                value=option,
                inline=False
            )

        view = MCQuestionView(self)
        await self.interaction.followup.send(embed=embed, view=view)

    async def handle_answer(self, answer: str):
        """Process the user's answer"""
        question = self.questions[self.current_question]
        correct_answer = question['correct']
        
        is_correct = answer == correct_answer
        if is_correct:
            self.score += 1
            response = "✅ Correct!"
        else:
            response = f"❌ Incorrect. The correct answer was {correct_answer}."

        await self.interaction.followup.send(response)
        
        self.current_question += 1
        await self.show_question()

    async def end_quiz(self):
        """End the quiz and show results"""
        percentage = (self.score / self.total_questions) * 100
        
        embed = discord.Embed(
            title="🎯 Quiz Complete!",
            description=f"You scored {self.score}/{self.total_questions} ({percentage:.1f}%)",
            color=discord.Color.green()
        )
        
        await self.interaction.followup.send(embed=embed)

async def setup(bot):
    await bot.add_cog(MCQuiz(bot)) 